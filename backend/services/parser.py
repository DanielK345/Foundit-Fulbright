import io
import os

import pdfplumber
import pypdfium2 as pdfium
from PIL import Image
from pptx import Presentation

try:
    import pytesseract
except ImportError:
    pytesseract = None


def _ocr_available() -> bool:
    if pytesseract is None:
        return False
    try:
        pytesseract.get_tesseract_version()
        return True
    except Exception:
        return False


def _clean_text_blocks(blocks: list[str]) -> list[str]:
    cleaned, seen = [], set()
    for block in blocks:
        text = block.strip()
        if not text:
            continue
        normalized = " ".join(text.split())
        if normalized in seen:
            continue
        seen.add(normalized)
        cleaned.append(text)
    return cleaned


def _ocr_image(image: Image.Image) -> str:
    if not _ocr_available():
        return ""
    grayscale = image.convert("L")
    return pytesseract.image_to_string(grayscale).strip()


def _render_pdf_page(page: pdfium.PdfPage, scale: float = 2.0) -> Image.Image:
    bitmap = page.render(scale=scale)
    try:
        return bitmap.to_pil()
    finally:
        bitmap.close()


def parse_pdf(file_path: str) -> list[dict]:
    """Extract text from each PDF page using pdfplumber, with Tesseract OCR fallback."""
    pages = []

    with pdfplumber.open(file_path) as pdf, pdfium.PdfDocument(file_path) as pdf_doc:
        for i, page in enumerate(pdf.pages):
            raw_text = (page.extract_text() or "").strip()

            blocks = []
            if raw_text:
                blocks.append(raw_text)

            # OCR fallback for image-heavy pages with little text
            if len(raw_text) < 40:
                pdf_page = pdf_doc[i]
                rendered = _render_pdf_page(pdf_page)
                try:
                    ocr_text = _ocr_image(rendered)
                finally:
                    rendered.close()
                    pdf_page.close()
                if ocr_text:
                    blocks.append(f"[OCR]\n{ocr_text}")

            content_blocks = _clean_text_blocks(blocks)
            content = "\n\n".join(content_blocks)
            extraction_method = "text" if raw_text else "ocr"

            if content:
                pages.append({
                    "source": f"page_{i + 1}",
                    "content": content,
                    "extraction_method": extraction_method,
                })

    return pages


def parse_pptx(file_path: str) -> list[dict]:
    """Extract text from each PPTX slide using python-pptx, with Tesseract OCR on embedded images."""
    prs = Presentation(file_path)
    slides = []

    for i, slide in enumerate(prs.slides):
        text_blocks = []
        embedded_images = []

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        text_blocks.append(text)

            if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                try:
                    img = Image.open(io.BytesIO(shape.image.blob))
                    embedded_images.append(img.copy())
                    img.close()
                except Exception:
                    pass

        notes_text = ""
        try:
            notes_slide = slide.notes_slide
            for para in notes_slide.notes_text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    notes_text += t + "\n"
            notes_text = notes_text.strip()
        except Exception:
            pass

        if not text_blocks and not embedded_images and not notes_text:
            continue

        blocks = list(text_blocks)
        has_ocr = False
        for img in embedded_images:
            ocr_text = _ocr_image(img)
            if ocr_text:
                blocks.append(f"[Image OCR]\n{ocr_text}")
                has_ocr = True
            img.close()

        if notes_text:
            blocks.append(f"[Speaker notes]\n{notes_text}")

        content_blocks = _clean_text_blocks(blocks)
        content = "\n\n".join(content_blocks)
        has_text = any(not b.startswith("[Image OCR]") for b in content_blocks)
        extraction_method = "text+ocr" if has_text and has_ocr else ("ocr" if has_ocr else "text")

        if content:
            slides.append({
                "source": f"slide_{i + 1}",
                "content": content,
                "extraction_method": extraction_method,
            })

    return slides


def parse_file(file_path: str, filename: str) -> list[dict]:
    """Route to the correct parser based on file extension."""
    ext = filename.lower().rsplit(".", 1)[-1]
    if ext == "pdf":
        return parse_pdf(file_path)
    elif ext == "pptx":
        return parse_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
